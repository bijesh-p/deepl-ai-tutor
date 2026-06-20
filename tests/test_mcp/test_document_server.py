"""Tests for document_server MCP tools via MCPClient.

Verifies each tool's output is a drop-in replacement for calling the
corresponding parser directly — i.e. Document.from_json() parses it with the
same title/sections/source_type.
"""
from __future__ import annotations

from pathlib import Path

from backend.core.mcp_client import get_client
from backend.ingestion.models import Document, SourceType

FIXTURE = Path(__file__).parent.parent / "fixtures" / "sample.pdf"
VTT_FIXTURE = Path(__file__).parent.parent / "fixtures" / "sample.vtt"


def test_extract_text_from_pdf_matches_parse_pdf():
    from backend.ingestion.pdf_parser import parse_pdf

    direct = parse_pdf(str(FIXTURE))

    doc_json = get_client("document_server").call(
        "extract_text_from_pdf", file_path=str(FIXTURE), max_pages=4
    )
    via_mcp = Document.from_json(doc_json)

    assert via_mcp.title == direct.title
    assert via_mcp.source_type == SourceType.PDF
    assert via_mcp.total_pages == direct.total_pages
    assert [s.title for s in via_mcp.sections] == [s.title for s in direct.sections]
    assert [s.body for s in via_mcp.sections] == [s.body for s in direct.sections]


def test_extract_text_from_pptx_matches_parse_pptx(tmp_path):
    from pptx import Presentation
    from backend.ingestion.pptx_parser import parse_pptx

    pptx_path = tmp_path / "sample.pptx"
    prs = Presentation()
    layout = prs.slide_layouts[1]  # title + content
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "Neural Networks"
    slide.placeholders[1].text = "A neural network is a computational model inspired by the brain."
    prs.save(str(pptx_path))

    direct = parse_pptx(str(pptx_path))
    doc_json = get_client("document_server").call(
        "extract_text_from_pptx", file_path=str(pptx_path)
    )
    via_mcp = Document.from_json(doc_json)

    assert via_mcp.source_type == SourceType.PPTX
    assert via_mcp.title == direct.title
    assert [s.title for s in via_mcp.sections] == [s.title for s in direct.sections]
    assert [s.body for s in via_mcp.sections] == [s.body for s in direct.sections]


def test_extract_text_from_docx_matches_parse_docx(tmp_path):
    import docx as _docx
    from backend.ingestion.docx_parser import parse_docx

    docx_path = tmp_path / "sample.docx"
    d = _docx.Document()
    d.add_heading("Introduction", level=1)
    d.add_paragraph("This section covers the foundational concepts of machine learning.")
    d.add_heading("Methods", level=1)
    d.add_paragraph("We describe the experimental setup and evaluation metrics here.")
    d.save(str(docx_path))

    direct = parse_docx(str(docx_path))
    doc_json = get_client("document_server").call(
        "extract_text_from_docx", file_path=str(docx_path)
    )
    via_mcp = Document.from_json(doc_json)

    assert via_mcp.source_type == SourceType.DOCX
    assert via_mcp.title == direct.title
    assert [s.title for s in via_mcp.sections] == [s.title for s in direct.sections]
    assert [s.body for s in via_mcp.sections] == [s.body for s in direct.sections]


def test_extract_text_from_vtt_matches_parse_vtt():
    from backend.ingestion.vtt_parser import parse_vtt

    direct = parse_vtt(str(VTT_FIXTURE))
    doc_json = get_client("document_server").call(
        "extract_text_from_vtt", file_path=str(VTT_FIXTURE)
    )
    via_mcp = Document.from_json(doc_json)

    assert via_mcp.source_type == SourceType.VTT
    assert via_mcp.total_pages == direct.total_pages
    assert [s.title for s in via_mcp.sections] == [s.title for s in direct.sections]
    assert [s.body for s in via_mcp.sections] == [s.body for s in direct.sections]
