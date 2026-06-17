from __future__ import annotations

import os
import tempfile

import pytest
from pptx import Presentation
from pptx.util import Inches

from backend.ingestion.models import SourceType
from backend.ingestion.pptx_parser import parse_pptx


def _make_pptx(slides: list[dict], title: str = "") -> str:
    """Create a minimal .pptx in a temp file and return its path.

    Each entry in *slides* is {"title": str, "body": str | None}.
    """
    prs = Presentation()
    if title:
        prs.core_properties.title = title

    blank_layout = prs.slide_layouts[5]  # "Blank" layout
    title_body_layout = prs.slide_layouts[1]  # "Title and Content"

    for slide_def in slides:
        slide = prs.slides.add_slide(title_body_layout)
        slide.shapes.title.text = slide_def.get("title", "")
        body = slide_def.get("body", "")
        if body:
            tf = slide.placeholders[1].text_frame
            tf.text = body

    fd, path = tempfile.mkstemp(suffix=".pptx")
    os.close(fd)
    prs.save(path)
    return path


def _cleanup(path: str) -> None:
    if os.path.exists(path):
        os.unlink(path)


class TestPptxParser:
    def test_basic_slides(self):
        path = _make_pptx(
            [
                {"title": "Intro", "body": "Welcome to the course."},
                {"title": "Topic A", "body": "Details about A."},
            ],
            title="My Presentation",
        )
        try:
            doc = parse_pptx(path)
            assert doc.title == "My Presentation"
            assert doc.source_type == SourceType.PPTX
            assert len(doc.sections) == 2
            assert doc.sections[0].title == "Intro"
            assert "Welcome" in doc.sections[0].body
            assert doc.sections[1].title == "Topic A"
            assert doc.total_pages == 2
        finally:
            _cleanup(path)

    def test_title_falls_back_to_filename_stem(self):
        path = _make_pptx([{"title": "Slide 1", "body": "Some text."}])
        try:
            doc = parse_pptx(path)
            # No core_properties.title set — should use stem of the temp filename
            assert doc.title  # not empty
            assert doc.source_filename.endswith(".pptx")
        finally:
            _cleanup(path)

    def test_slide_title_falls_back_to_slide_number(self):
        # Build a slide without a title placeholder by using the blank layout
        prs = Presentation()
        blank = prs.slide_layouts[5]
        slide = prs.slides.add_slide(blank)
        # Add a plain text box as body (no title shape)
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        txBox.text_frame.text = "Body text only"

        fd, path = tempfile.mkstemp(suffix=".pptx")
        os.close(fd)
        prs.save(path)
        try:
            doc = parse_pptx(path)
            assert len(doc.sections) == 1
            assert doc.sections[0].title == "Slide 1"
        finally:
            _cleanup(path)

    def test_empty_body_slides_are_skipped(self):
        path = _make_pptx(
            [
                {"title": "Title Only", "body": ""},
                {"title": "Has Content", "body": "Real content here."},
            ]
        )
        try:
            doc = parse_pptx(path)
            # Slide with empty body should be skipped
            assert len(doc.sections) == 1
            assert doc.sections[0].title == "Has Content"
        finally:
            _cleanup(path)

    def test_max_slides_cap(self):
        slides = [{"title": f"Slide {i}", "body": f"Content {i}."} for i in range(1, 6)]
        path = _make_pptx(slides)
        try:
            doc = parse_pptx(path, max_slides=3)
            assert len(doc.sections) <= 3
        finally:
            _cleanup(path)

    def test_source_type_and_filename(self):
        path = _make_pptx([{"title": "A", "body": "B."}])
        try:
            doc = parse_pptx(path)
            assert doc.source_type == SourceType.PPTX
            assert doc.source_filename.endswith(".pptx")
        finally:
            _cleanup(path)

    def test_document_roundtrip_json(self):
        from backend.ingestion.models import Document
        path = _make_pptx([{"title": "JSON Test", "body": "Roundtrip body."}], title="RT")
        try:
            doc = parse_pptx(path)
            recovered = Document.from_json(doc.to_json())
            assert recovered.title == doc.title
            assert recovered.source_type == SourceType.PPTX
            assert len(recovered.sections) == len(doc.sections)
        finally:
            _cleanup(path)
