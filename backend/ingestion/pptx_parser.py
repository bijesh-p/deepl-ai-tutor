from __future__ import annotations

import itertools
import uuid
from pathlib import Path

from pptx import Presentation

from backend.ingestion.models import Document, Section, SourceType

_DEFAULT_MAX_SLIDES = 16


def parse_pptx(file_path: str, max_slides: int = _DEFAULT_MAX_SLIDES) -> Document:
    """Parse a PowerPoint file into the unified Document model.

    Each slide becomes a Section. Slides with no extractable text body are
    skipped. At most *max_slides* slides are processed.
    """
    path = Path(file_path)
    doc_id = str(uuid.uuid4())

    prs = Presentation(str(path))

    raw_title = (prs.core_properties.title or "").strip()
    title = raw_title if raw_title else path.stem

    sections: list[Section] = []
    for slide_num, slide in enumerate(itertools.islice(prs.slides, max_slides), start=1):
        slide_title = _slide_title(slide, slide_num)
        body = _slide_body(slide)
        if not body:
            continue
        sections.append(
            Section(
                section_id=str(uuid.uuid4()),
                title=slide_title,
                body=body,
                level=1,
            )
        )

    return Document(
        doc_id=doc_id,
        title=title,
        source_filename=path.name,
        source_type=SourceType.PPTX,
        sections=sections,
        total_pages=len(sections),
    )


def _slide_title(slide, slide_num: int) -> str:
    shape = slide.shapes.title
    if shape is not None:
        text = shape.text.strip()
        if text:
            return text
    return f"Slide {slide_num}"


def _slide_body(slide) -> str:
    parts: list[str] = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if shape.is_placeholder and shape.placeholder_format.idx == 0:
            continue  # title placeholder — already captured by _slide_title
        text = shape.text_frame.text.strip()
        if text:
            parts.append(text)
    return "\n\n".join(parts)
