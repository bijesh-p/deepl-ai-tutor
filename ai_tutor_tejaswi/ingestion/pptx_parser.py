from __future__ import annotations
import os
import uuid

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from ingestion.image_extractor import save_image
from ingestion.models import Document, Section, SourceType


def parse_pptx(file_path: str) -> Document:
    doc_id = str(uuid.uuid4())
    upload_dir = os.path.join(
        os.environ.get("AI_TUTOR_UPLOAD_DIR", "data/uploads"), doc_id, "images"
    )

    prs = Presentation(file_path)
    sections: list[Section] = []

    for slide_num, slide in enumerate(prs.slides):
        title_text = ""
        body_parts: list[str] = []
        images = []

        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    ext_img = save_image(
                        shape.image.blob, upload_dir, None, f"slide {slide_num + 1}"
                    )
                    images.append(ext_img)
                except Exception:
                    pass
                continue

            if not shape.has_text_frame:
                continue

            ph = getattr(shape, "placeholder_format", None)
            text = shape.text_frame.text.strip()
            if not text:
                continue

            if ph is not None and ph.idx == 0:
                title_text = text
            else:
                body_parts.append(text)

        notes = ""
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()

        sections.append(
            Section(
                section_id=str(uuid.uuid4()),
                title=title_text or f"Slide {slide_num + 1}",
                body="\n".join(body_parts),
                level=1,
                images=images,
                metadata={"speaker_notes": notes} if notes else {},
            )
        )

    doc_title = os.path.splitext(os.path.basename(file_path))[0]
    return Document(
        doc_id=doc_id,
        title=doc_title,
        source_filename=os.path.basename(file_path),
        source_type=SourceType.PPTX,
        sections=sections,
        total_pages=len(prs.slides),
    )
